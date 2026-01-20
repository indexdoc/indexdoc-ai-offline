from typing import List

from domain.kb_domain.serv.DocLoad.DocLoadImp.BaseLoader import BaseLoader


class OCRPPTLoader(BaseLoader):
    def __init__(self, file_path: str):
        super().__init__(file_path)
        self.unstructured_kwargs = None

    def _load_impl(self) -> List:
        def ppt2text(filepath):
            # pip install python-pptx
            from pptx import Presentation
            from PIL import Image
            import numpy as np
            from io import BytesIO
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                sorted_shapes = sorted(slide.shapes,
                                       key=lambda x: (x.top, x.left))  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
            return resp

        text = ppt2text(self.file_path)
        return [text]
