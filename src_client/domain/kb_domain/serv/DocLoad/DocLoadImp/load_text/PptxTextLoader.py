from typing import List
from pptx import Presentation
from domain.kb_domain.serv.DocLoad.DocLoadImp.BaseLoader import BaseLoader


class PptxTextLoader(BaseLoader):
    def __init__(self, file_path: str):
        super().__init__(file_path)

    def _load_impl(self) -> List[str]:
        def ppt2text(filepath: str) -> str:
            prs = Presentation(filepath)
            text = ""

            def extract_text(shape):
                nonlocal text
                # 文本框
                if shape.has_text_frame:
                    text += shape.text.strip() + "\n"

                # 表格
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text += cell.text.strip() + "\n"

                # 组合形状（递归）
                if hasattr(shape, "shapes"):
                    for child in shape.shapes:
                        extract_text(child)

            # 遍历每一页幻灯片
            for slide in prs.slides:
                sorted_shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
                for shape in sorted_shapes:
                    extract_text(shape)

            return text

        return [ppt2text(self.file_path)]


